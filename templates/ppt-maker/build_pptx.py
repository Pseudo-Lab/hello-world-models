"""DreamerV2 강의 .md를 review.css 스타일에 맞춘 편집 가능 PPTX로 변환."""
import hashlib
import re
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import mathtext
from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from PIL import Image


def set_run_font(run, name: str):
    """Run에 Latin/EA/CS 폰트를 모두 설정 (한글까지 적용되도록)."""
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        full_tag = qn(f"a:{tag}")
        elem = rPr.find(full_tag)
        if elem is None:
            elem = etree.SubElement(rPr, full_tag)
        elem.set("typeface", name)


def _get_or_create_pPr(paragraph):
    p = paragraph._p
    pPr = p.find(qn("a:pPr"))
    if pPr is None:
        pPr = etree.SubElement(p, qn("a:pPr"))
        p.remove(pPr)
        p.insert(0, pPr)
    return pPr


def set_bullet(paragraph, level: int = 0):
    """OOXML 정식 bullet 설정 — Arial을 bullet 폰트로 지정해 글리프 누락 방지."""
    pPr = _get_or_create_pPr(paragraph)
    for tag in ("buFont", "buChar", "buNone", "buAutoNum"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    if level > 0:
        pPr.set("lvl", str(level))
    pPr.set("marL", str(285750 * (level + 1)))
    pPr.set("indent", "-285750")
    bullet_chars = ["•", "–", "·"]
    ch = bullet_chars[min(level, 2)]
    buFont = etree.SubElement(pPr, qn("a:buFont"))
    buFont.set("typeface", "Arial")
    buChar = etree.SubElement(pPr, qn("a:buChar"))
    buChar.set("char", ch)


def set_no_bullet(paragraph):
    """Paragraph에서 bullet 제거."""
    pPr = _get_or_create_pPr(paragraph)
    for tag in ("buFont", "buChar", "buAutoNum"):
        for el in pPr.findall(qn(f"a:{tag}")):
            pPr.remove(el)
    if pPr.find(qn("a:buNone")) is None:
        etree.SubElement(pPr, qn("a:buNone"))


# ===== 스타일 =====
WINE = RGBColor(0x8A, 0x15, 0x38)
WINE_LIGHT = RGBColor(0xA8, 0x20, 0x50)
WINE_BG = RGBColor(0xFD, 0xF2, 0xF5)
DARK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x55, 0x55, 0x55)
GRAY_LIGHT = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_MAIN = "Freesentation"  # 영문명 사용 (한글 family명 "프리젠테이션"과 동일 폰트)
FONT_FALLBACK = "맑은 고딕"
FONT_CODE = "JetBrains Mono"

# 16:9 1280x720 비율 (Marp 기본)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

TITLE_LEFT = Inches(0.55)
TITLE_TOP = Inches(0.4)
TITLE_W = Inches(12.2)
TITLE_H = Inches(0.55)

CONTENT_LEFT = Inches(0.6)
CONTENT_TOP = Inches(1.2)
CONTENT_W = Inches(12.1)
CONTENT_H = Inches(5.9)

ROOT = Path(__file__).parent
MD_PATH = ROOT / "DreamerV2_lecture.md"
OUT_PATH = ROOT / "DreamerV2_lecture_v2.pptx"
EQ_DIR = ROOT / "visual" / "equations"


def render_latex_to_png(latex: str, *, fontsize: int = 22, color: str = "#111111",
                        dpi: int = 300) -> Path:
    """LaTeX 수식 → PNG 이미지 (캐싱)."""
    EQ_DIR.mkdir(parents=True, exist_ok=True)
    # 캐시 키
    h = hashlib.md5(f"{latex}|{fontsize}|{color}|{dpi}".encode()).hexdigest()[:12]
    out_path = EQ_DIR / f"eq_{h}.png"
    if out_path.exists():
        return out_path

    # matplotlib mathtext로 렌더링
    # \dot, \begin{cases} 등 일부 LaTeX 명령은 mathtext에서 미지원
    # → 호환되지 않는 부분은 사전 치환
    expr = latex.strip()
    expr = expr.replace("\\doteq", "\\equiv")  # mathtext 미지원 치환
    expr = expr.replace("\\big[", "[").replace("\\big]", "]")
    expr = expr.replace("\\Big[", "[").replace("\\Big]", "]")

    # \begin{cases} ... \end{cases} 처리 (mathtext 미지원)
    cases_match = re.search(r"\\begin\{cases\}(.+?)\\end\{cases\}", expr, re.S)
    if cases_match:
        inner = cases_match.group(1)
        # &로 정렬된 분기를 줄바꿈으로 묶음 처리
        lines = [ln.strip() for ln in inner.split("\\\\") if ln.strip()]
        joined = "\\,\\,".join(ln.replace("&", "\\,\\,") for ln in lines)
        expr = expr[:cases_match.start()] + "\\{ " + joined + " \\}" + expr[cases_match.end():]

    fig = plt.figure(figsize=(0.01, 0.01))  # 임시 사이즈
    fig.patch.set_alpha(0)
    try:
        text_obj = fig.text(0, 0, f"${expr}$", fontsize=fontsize, color=color)
        fig.canvas.draw()
        # 텍스트 bbox로 figure 크기 재조정
        bbox = text_obj.get_window_extent(renderer=fig.canvas.get_renderer())
        w_in = bbox.width / fig.dpi + 0.2
        h_in = bbox.height / fig.dpi + 0.1
        plt.close(fig)
        fig = plt.figure(figsize=(w_in, h_in), dpi=dpi)
        fig.patch.set_alpha(0)
        fig.text(0.05, 0.5, f"${expr}$", fontsize=fontsize, color=color,
                 va="center", ha="left")
        fig.savefig(out_path, dpi=dpi, transparent=True, bbox_inches="tight",
                    pad_inches=0.05)
        plt.close(fig)
    except (mathtext.MathTextWarning, ValueError, Exception) as e:
        plt.close("all")
        # 렌더링 실패 시 일반 텍스트 fallback
        fig = plt.figure(figsize=(8, 1), dpi=dpi)
        fig.patch.set_alpha(0)
        fig.text(0.5, 0.5, latex, fontsize=14, color=color, ha="center", va="center",
                 family="monospace")
        fig.savefig(out_path, dpi=dpi, transparent=True, bbox_inches="tight",
                    pad_inches=0.05)
        plt.close(fig)
    return out_path


# ===== 마크다운 파싱 =====
@dataclass
class Block:
    kind: str  # 'bullet', 'para', 'image', 'table', 'code', 'math', 'quote'
    data: object = None


@dataclass
class Slide:
    title: str = ""
    subtitle: str = ""
    is_title_slide: bool = False
    is_toc: bool = False
    is_summary: bool = False
    blocks: List[Block] = field(default_factory=list)


def parse_markdown(text: str) -> List[Slide]:
    # 프론트매터 제거
    text = re.sub(r"^---\nmarp:.*?\n---\n", "", text, count=1, flags=re.S)
    raw_slides = [s.strip() for s in re.split(r"\n---\n", text) if s.strip()]
    slides = []
    for idx, raw in enumerate(raw_slides):
        slide = Slide()
        lines = raw.split("\n")
        i = 0
        # title
        while i < len(lines) and not lines[i].startswith("# "):
            i += 1
        if i < len(lines):
            slide.title = lines[i][2:].strip()
            i += 1
        # subtitle (h2)
        while i < len(lines) and lines[i].strip() == "":
            i += 1
        if i < len(lines) and lines[i].startswith("## "):
            slide.subtitle = lines[i][3:].strip()
            i += 1

        # 첫 슬라이드는 타이틀 슬라이드
        if idx == 0:
            slide.is_title_slide = True
        if slide.title == "목차":
            slide.is_toc = True
        if slide.title == "요약" or slide.title.startswith("한눈에 보기"):
            slide.is_summary = True

        # blocks
        bullet_buf = []
        para_buf = []
        table_buf = []
        code_buf = []
        in_code = False
        in_math = False
        math_buf = []

        def flush_bullets():
            if bullet_buf:
                slide.blocks.append(Block("bullet", list(bullet_buf)))
                bullet_buf.clear()

        def flush_paras():
            # 각 라인을 별도 para Block으로 (목차 등에서 라인별 처리 필요)
            for line in para_buf:
                slide.blocks.append(Block("para", line))
            para_buf.clear()

        def flush_table():
            if table_buf:
                slide.blocks.append(Block("table", list(table_buf)))
                table_buf.clear()

        while i < len(lines):
            line = lines[i]
            stripped = line.strip()

            if stripped.startswith("<!--"):
                i += 1
                continue
            if stripped.startswith("```"):
                if in_code:
                    flush_bullets(); flush_paras(); flush_table()
                    slide.blocks.append(Block("code", "\n".join(code_buf)))
                    code_buf = []
                    in_code = False
                else:
                    flush_bullets(); flush_paras(); flush_table()
                    in_code = True
                i += 1
                continue
            if in_code:
                code_buf.append(line)
                i += 1
                continue

            if stripped.startswith("$$"):
                if in_math:
                    flush_bullets(); flush_paras(); flush_table()
                    slide.blocks.append(Block("math", "\n".join(math_buf)))
                    math_buf = []
                    in_math = False
                else:
                    flush_bullets(); flush_paras(); flush_table()
                    in_math = True
                i += 1
                continue
            if in_math:
                math_buf.append(line)
                i += 1
                continue

            # 이미지 ![alt](path)
            m_img = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", stripped)
            if m_img:
                flush_bullets(); flush_paras(); flush_table()
                slide.blocks.append(Block("image", (m_img.group(1), m_img.group(2))))
                i += 1
                continue

            # 인용문
            if stripped.startswith("> "):
                flush_bullets(); flush_paras(); flush_table()
                quote_lines = []
                while i < len(lines) and lines[i].strip().startswith("> "):
                    quote_lines.append(lines[i].strip()[2:])
                    i += 1
                slide.blocks.append(Block("quote", " ".join(quote_lines)))
                continue

            # 테이블
            if stripped.startswith("|") and "|" in stripped[1:]:
                flush_bullets(); flush_paras()
                table_buf.append(stripped)
                i += 1
                continue
            else:
                if table_buf:
                    flush_table()

            # 불릿
            m_bullet = re.match(r"^(\s*)[-*]\s+(.*)$", line)
            if m_bullet:
                flush_paras()
                indent = len(m_bullet.group(1)) // 2
                bullet_buf.append((indent, m_bullet.group(2)))
                i += 1
                continue

            # 번호 리스트 (TOC 포함)
            m_num = re.match(r"^(\s*)(\d+)\.\s+(.*)$", line)
            if m_num:
                flush_paras()
                indent = len(m_num.group(1)) // 2
                bullet_buf.append((indent, f"{m_num.group(2)}. {m_num.group(3)}"))
                i += 1
                continue

            # 일반 문단
            if stripped:
                flush_bullets()
                para_buf.append(stripped)
            else:
                flush_bullets(); flush_paras()
            i += 1

        flush_bullets(); flush_paras(); flush_table()
        slides.append(slide)
    return slides


# ===== 텍스트 inline 파싱 (**bold**) =====
def add_runs_with_bold(paragraph, text: str, *, base_size=18, base_color=DARK,
                       bold_color=None, font=FONT_MAIN):
    """**...** 패턴을 파싱해 bold runs로 추가."""
    parts = re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text)
    bold_color = bold_color or WINE
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run()
            run.text = part[2:-2]
            set_run_font(run, font)
            run.font.size = Pt(base_size)
            run.font.bold = True
            run.font.color.rgb = bold_color
        elif part.startswith("`") and part.endswith("`"):
            run = paragraph.add_run()
            run.text = part[1:-1]
            set_run_font(run, FONT_CODE)
            run.font.size = Pt(base_size - 1)
            run.font.color.rgb = WINE
        else:
            run = paragraph.add_run()
            run.text = part
            set_run_font(run, font)
            run.font.size = Pt(base_size)
            run.font.color.rgb = base_color


# ===== 슬라이드 빌더 =====
def add_title(slide_obj, title_text: str):
    # 제목
    box = slide_obj.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_W, TITLE_H)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    set_run_font(run, FONT_MAIN)
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WINE

    # 와인 밑줄 (사각형 4px) — 제목 바로 아래
    underline = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        TITLE_LEFT, TITLE_TOP + TITLE_H,
        TITLE_W, Emu(45000),
    )
    underline.fill.solid()
    underline.fill.fore_color.rgb = WINE
    underline.line.fill.background()
    underline.shadow.inherit = False


def add_bullet_textbox(slide_obj, bullets: List[Tuple[int, str]], top: Emu = None,
                       height: Emu = None, base_size: int = 18):
    top = top if top is not None else CONTENT_TOP
    height = height if height is not None else CONTENT_H
    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    for idx, (indent, text) in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        size = base_size if indent == 0 else max(base_size - 2 * indent, 11)
        p.space_after = Pt(4)
        # 번호 리스트는 bullet 없이
        if re.match(r"^\d+\.\s", text):
            set_no_bullet(p)
        else:
            set_bullet(p, level=indent)
        add_runs_with_bold(p, text, base_size=size)
    return box


def resolve_image_path(path: str) -> Optional[Path]:
    """visual/01_X.png → v2/v2_01_X.png 또는 v1/01_X.png 우선순위로 검색."""
    direct = ROOT / path
    if direct.exists():
        return direct
    fname = Path(path).name  # 01_DreamerV2_전체_구조_개요.png
    # v2: v2_01_X.png 또는 extra_v2_01_X.png
    if fname.startswith("extra_"):
        v2_name = fname.replace("extra_", "extra_v2_", 1)
    else:
        v2_name = "v2_" + fname
    v2_path = ROOT / "visual" / "v2" / v2_name
    if v2_path.exists():
        return v2_path
    v1_path = ROOT / "visual" / "v1" / fname
    if v1_path.exists():
        return v1_path
    return None


def add_image(slide_obj, alt: str, path: str, top: Emu = None):
    """이미지를 alt(w80 등)에 따라 크기 조정해서 중앙 배치."""
    img_path = resolve_image_path(path)
    if img_path is None:
        print(f"  [WARN] image not found: {path}")
        return
    # alt에서 width 추출
    m = re.search(r"w(\d+)", alt)
    width_pct = int(m.group(1)) if m else 80
    max_w = SLIDE_W * (width_pct / 100.0)

    with Image.open(img_path) as im:
        ratio = im.height / im.width
    img_w = max_w
    img_h = int(img_w * ratio)

    # 사용 가능한 영역
    avail_top = top if top is not None else CONTENT_TOP
    avail_h = SLIDE_H - avail_top - Inches(0.4)
    if img_h > avail_h:
        img_h = int(avail_h)
        img_w = int(img_h / ratio)

    left = int((SLIDE_W - img_w) / 2)
    top_pos = int(avail_top + (avail_h - img_h) / 2)
    slide_obj.shapes.add_picture(str(img_path), left, top_pos, img_w, img_h)


def add_table(slide_obj, table_lines: List[str]):
    # Header | --- | rows
    rows = []
    for line in table_lines:
        cells = [c.strip() for c in line.strip("|").split("|")]
        rows.append(cells)
    # 두번째 줄(--- 정렬)은 제외
    if len(rows) >= 2 and all(re.match(r"^:?-+:?$", c.replace(" ", "")) for c in rows[1] if c):
        header = rows[0]
        data_rows = rows[2:]
    else:
        header = rows[0]
        data_rows = rows[1:]

    n_cols = len(header)
    n_rows = 1 + len(data_rows)

    # 사용 가능 폭에서 균등 분배
    table_w = CONTENT_W
    table_h = min(CONTENT_H, Inches(0.5 * n_rows + 0.3))
    table_top = CONTENT_TOP + (CONTENT_H - table_h) // 2

    table_shape = slide_obj.shapes.add_table(
        n_rows, n_cols, CONTENT_LEFT, table_top, table_w, table_h
    )
    tbl = table_shape.table

    # 헤더 스타일
    for ci, txt in enumerate(header):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = WINE
        tf = cell.text_frame
        tf.margin_left = Emu(50000); tf.margin_right = Emu(50000)
        tf.margin_top = Emu(30000); tf.margin_bottom = Emu(30000)
        tf.text = ""
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_runs_with_bold(p, txt, base_size=14, base_color=WHITE, bold_color=WHITE)

    # 데이터
    for ri, row in enumerate(data_rows, start=1):
        for ci in range(n_cols):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xFA, 0xFA, 0xFA)
            tf = cell.text_frame
            tf.margin_left = Emu(50000); tf.margin_right = Emu(50000)
            tf.margin_top = Emu(20000); tf.margin_bottom = Emu(20000)
            tf.text = ""
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            txt = row[ci] if ci < len(row) else ""
            add_runs_with_bold(p, txt, base_size=12, base_color=DARK)


def add_paragraph_box(slide_obj, text: str, top: Emu, height: Emu, *,
                      align=PP_ALIGN.LEFT, base_size=18, italic=False, bg=None,
                      border=None):
    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, height)
    if bg:
        box.fill.solid()
        box.fill.fore_color.rgb = bg
    if border:
        box.line.color.rgb = border
        box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(150000); tf.margin_right = Emu(150000)
    tf.margin_top = Emu(100000); tf.margin_bottom = Emu(100000)
    p = tf.paragraphs[0]
    p.alignment = align
    add_runs_with_bold(p, text, base_size=base_size)
    if italic:
        for run in p.runs:
            run.font.italic = True
    return box


def add_quote(slide_obj, text: str):
    # 중앙 인용 박스
    h = Inches(2.2)
    top = CONTENT_TOP + (CONTENT_H - h) // 2
    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
    tf.margin_top = Inches(0.3); tf.margin_bottom = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_runs_with_bold(p, text, base_size=24, base_color=WINE, bold_color=WINE)
    for run in p.runs:
        run.font.italic = True
    # 좌측 와인색 막대
    bar = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        CONTENT_LEFT - Emu(80000), top, Emu(60000), h
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WINE
    bar.line.fill.background()
    bar.shadow.inherit = False


def add_code(slide_obj, code: str, top: Emu, height: Emu):
    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    box.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2); tf.margin_bottom = Inches(0.2)
    lines = code.split("\n")
    for idx, ln in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln if ln else " "
        set_run_font(run, FONT_CODE)
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0xC7, 0x25, 0x4E)


# ===== 슬라이드 타입별 빌더 =====
def build_title_slide(slide_obj, slide: Slide):
    # 와인 사이드 바
    bar = slide_obj.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.45), SLIDE_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = WINE
    bar.line.fill.background()
    bar.shadow.inherit = False

    # 큰 제목 (2줄까지 허용, 폰트 약간 축소)
    box = slide_obj.shapes.add_textbox(
        Inches(1.2), Inches(2.0), Inches(11.5), Inches(2.0)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = slide.title
    set_run_font(run, FONT_MAIN)
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = WINE

    # 부제목 (제목 아래)
    if slide.subtitle:
        box2 = slide_obj.shapes.add_textbox(
            Inches(1.2), Inches(4.2), Inches(11.5), Inches(0.7)
        )
        tf2 = box2.text_frame
        tf2.margin_left = Emu(0); tf2.margin_right = Emu(0); tf2.margin_top = Emu(0)
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = slide.subtitle
        set_run_font(run2, FONT_MAIN)
        run2.font.size = Pt(24)
        run2.font.color.rgb = DARK

    # 본문(저자 등) - paragraph blocks 결합
    para_lines = [b.data for b in slide.blocks if b.kind == "para"]
    if para_lines:
        box3 = slide_obj.shapes.add_textbox(
            Inches(1.2), Inches(5.4), Inches(11.5), Inches(1.5)
        )
        tf3 = box3.text_frame
        tf3.word_wrap = True
        tf3.margin_left = Emu(0); tf3.margin_right = Emu(0); tf3.margin_top = Emu(0)
        for idx, line in enumerate(para_lines):
            p3 = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
            p3.space_after = Pt(2)
            add_runs_with_bold(p3, line, base_size=15, base_color=GRAY)


def estimate_block_lines(blk: Block) -> int:
    if blk.kind == "para":
        return 1
    if blk.kind == "bullet":
        return len(blk.data)
    if blk.kind == "code":
        return blk.data.count("\n") + 1
    return 1


def render_text_blocks(slide_obj, blocks: List[Block], top, height,
                        base_size: int = 18):
    """여러 텍스트 블록(para+bullet)을 하나의 텍스트박스에 순서대로 렌더링."""
    if not blocks:
        return
    total_lines = sum(estimate_block_lines(b) for b in blocks)
    if base_size >= 18 and total_lines > 9:
        base_size = 16
    if base_size >= 16 and total_lines > 14:
        base_size = 14
    if base_size >= 14 and total_lines > 20:
        base_size = 13

    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, top, CONTENT_W, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    first = True
    for blk in blocks:
        if blk.kind == "para":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            was_first = first
            first = False
            p.space_before = Pt(0 if was_first else 6)
            p.space_after = Pt(4)
            set_no_bullet(p)
            add_runs_with_bold(p, blk.data, base_size=base_size)
        elif blk.kind == "bullet":
            for indent, text in blk.data:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                sz = base_size if indent == 0 else max(base_size - 2 * indent, 11)
                p.space_after = Pt(3)
                if re.match(r"^\d+\.\s", text):
                    set_no_bullet(p)
                else:
                    set_bullet(p, level=indent)
                add_runs_with_bold(p, text, base_size=sz)


def build_content_slide(slide_obj, slide: Slide):
    add_title(slide_obj, slide.title)

    blocks = slide.blocks
    images = [b for b in blocks if b.kind == "image"]
    tables = [b for b in blocks if b.kind == "table"]
    quotes = [b for b in blocks if b.kind == "quote"]
    codes = [b for b in blocks if b.kind == "code"]
    maths = [b for b in blocks if b.kind == "math"]

    # 1) 이미지 단독 (또는 짧은 인트로 + 이미지)
    text_blocks = [b for b in blocks if b.kind in ("para", "bullet")]
    if images and not (tables or quotes or codes or maths) and not any(b.kind == "bullet" for b in text_blocks):
        intro_text = [b.data for b in text_blocks]
        img_top = CONTENT_TOP
        if intro_text:
            box = slide_obj.shapes.add_textbox(
                CONTENT_LEFT, CONTENT_TOP, CONTENT_W, Inches(0.55)
            )
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0)
            for idx, line in enumerate(intro_text):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                add_runs_with_bold(p, line, base_size=15, base_color=GRAY)
            img_top = CONTENT_TOP + Inches(0.7)
        alt, path = images[0].data
        add_image(slide_obj, alt, path, top=img_top)
        return

    # 2) 테이블
    if tables:
        add_table(slide_obj, tables[0].data)
        tail = [b for b in blocks if b.kind in ("bullet", "para")]
        if tail:
            tail_top = CONTENT_TOP + Inches(4.6)
            render_text_blocks(slide_obj, tail, top=tail_top, height=Inches(1.2), base_size=13)
        return

    # 3) 인용문
    if quotes:
        add_quote(slide_obj, quotes[0].data)
        tail = [b for b in blocks if b.kind in ("bullet", "para")]
        if tail:
            render_text_blocks(slide_obj, tail, top=CONTENT_TOP + Inches(3.8),
                               height=Inches(1.7), base_size=15)
        return

    # 4) 코드 (Straight-Through 슬라이드)
    if codes:
        # 코드 위쪽에 첫 텍스트 블록, 아래에 나머지
        code_idx = blocks.index(codes[0])
        above = [b for b in blocks[:code_idx] if b.kind in ("para", "bullet")]
        below = [b for b in blocks[code_idx + 1:] if b.kind in ("para", "bullet")]
        if above:
            render_text_blocks(slide_obj, above, top=CONTENT_TOP, height=Inches(2.0), base_size=15)
        code_top = CONTENT_TOP + (Inches(2.1) if above else Inches(0.2))
        code_h = Inches(1.6)
        add_code(slide_obj, codes[0].data, code_top, code_h)
        if below:
            render_text_blocks(slide_obj, below, top=code_top + code_h + Inches(0.2),
                               height=Inches(1.7), base_size=14)
        return

    # 5) 수식 (λ-return)
    if maths:
        math_idx = blocks.index(maths[0])
        above = [b for b in blocks[:math_idx] if b.kind in ("para", "bullet")]
        below = [b for b in blocks[math_idx + 1:] if b.kind in ("para", "bullet")]
        if above:
            render_text_blocks(slide_obj, above, top=CONTENT_TOP, height=Inches(2.2), base_size=15)
        math_top = CONTENT_TOP + (Inches(2.4) if above else Inches(0.4))
        math_h = Inches(1.4)
        # 배경 박스
        math_bg = slide_obj.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, CONTENT_LEFT, math_top, CONTENT_W, math_h
        )
        math_bg.fill.solid()
        math_bg.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        math_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        math_bg.line.width = Pt(0.75)
        math_bg.shadow.inherit = False
        # LaTeX → PNG 후 중앙 배치
        eq_png = render_latex_to_png(maths[0].data, fontsize=20, color="#8a1538")
        with Image.open(eq_png) as im:
            ratio = im.height / im.width
        max_eq_w = CONTENT_W - Inches(1.0)
        eq_w = min(max_eq_w, Inches(8.5))
        eq_h = int(eq_w * ratio)
        if eq_h > math_h - Inches(0.3):
            eq_h = math_h - Inches(0.3)
            eq_w = int(eq_h / ratio)
        eq_left = int(CONTENT_LEFT + (CONTENT_W - eq_w) / 2)
        eq_top = int(math_top + (math_h - eq_h) / 2)
        slide_obj.shapes.add_picture(str(eq_png), eq_left, eq_top, eq_w, eq_h)
        if below:
            render_text_blocks(slide_obj, below, top=math_top + math_h + Inches(0.2),
                               height=Inches(1.6), base_size=14)
        return

    # 6) 일반 텍스트 (para + bullet 혼합)
    if text_blocks:
        render_text_blocks(slide_obj, text_blocks, top=CONTENT_TOP, height=CONTENT_H)


def build_toc_slide(slide_obj, slide: Slide):
    add_title(slide_obj, slide.title)
    # TOC 항목은 paragraph 또는 bullet으로 들어옴
    items = []
    for b in slide.blocks:
        if b.kind == "para":
            items.append(b.data)
        elif b.kind == "bullet":
            for indent, text in b.data:
                items.append(text)
    if not items:
        return
    box = slide_obj.shapes.add_textbox(CONTENT_LEFT, CONTENT_TOP, CONTENT_W, CONTENT_H)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, text in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        add_runs_with_bold(p, text, base_size=20, base_color=DARK, bold_color=WINE)


# ===== 메인 =====
def main():
    md_text = MD_PATH.read_text(encoding="utf-8")
    slides = parse_markdown(md_text)
    print(f"파싱: {len(slides)}개 슬라이드")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    blank_layout = prs.slide_layouts[6]  # 빈 레이아웃

    for idx, slide in enumerate(slides):
        s = prs.slides.add_slide(blank_layout)
        if slide.is_title_slide:
            build_title_slide(s, slide)
        elif slide.is_toc:
            build_toc_slide(s, slide)
        else:
            build_content_slide(s, slide)
        print(f"  [{idx+1}/{len(slides)}] {slide.title[:40]}")

    prs.save(str(OUT_PATH))
    print(f"저장: {OUT_PATH}")


if __name__ == "__main__":
    main()
